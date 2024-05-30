import sys
from PyQt5 import QtGui
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QTextEdit, QFileDialog, QLabel, QHBoxLayout, QDialog, QTableWidget, QTableWidgetItem, QStackedWidget, QMessageBox, QApplication,QDesktopWidget
from PyQt5.QtGui import QIcon
from PyQt5.QtCore import QUrl
import os
import glob
import time
import olefile
import zlib
import struct
import joblib
import re
from pptx import Presentation
from PIL import Image
import pytesseract
from docx import Document
from xml.etree.ElementTree import parse
from pdfminer.high_level import extract_text
import zipfile
from konlpy.tag import Okt


def read_pptx(file_path):
    pres = Presentation(file_path)
    text = ''
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    text = re.sub('[^가-힣a-zA-Z0-9\s]', '', text)
    text = text.replace('\x0b', '')
    return text

def read_docx(file_path):
    doc = Document(file_path)
    fullText = []

    # Process paragraphs
    for para in doc.paragraphs:
        clean_para = re.sub('[^가-힣a-zA-Z0-9\s.,()-:<>]', '', para.text)
        fullText.append(clean_para)

    # Process tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    clean_cell_text = re.sub('[^가-힣a-zA-Z0-9\s.,()-:<>]', '', paragraph.text)
                    fullText.append(clean_cell_text)

    return '\n'.join(fullText)

def ocr_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image, lang='kor')
    return text

def get_hwp_text(filename):
    f = olefile.OleFileIO(filename)
    dirs = f.listdir()

    if ["FileHeader"] not in dirs or \
       ["\x05HwpSummaryInformation"] not in dirs:
        raise Exception("Not Valid HWP.")

    header = f.openstream("FileHeader")
    header_data = header.read()
    is_compressed = (header_data[36] & 1) == 1

    nums = []
    for d in dirs:
        if d[0] == "BodyText":
           nums.append(int(d[1][len("Section"):]))
    sections = ["BodyText/Section"+str(x) for x in sorted(nums)]

    text = ""
    for section in sections:
        bodytext = f.openstream(section)
        data = bodytext.read()
        if is_compressed:
           unpacked_data = zlib.decompress(data, -15)
        else:
           unpacked_data = data

        section_text = ""
        i = 0
        size = len(unpacked_data)
        while i < size:
            header = struct.unpack_from("<I", unpacked_data, i)[0]
            rec_type = header & 0x3ff
            rec_len = (header >> 20) & 0xfff
            if rec_type in [67]:
                rec_data = unpacked_data[i+4:i+4+rec_len]
                try:
                    section_text += rec_data.decode('utf-16')
                except UnicodeDecodeError:
                    pass
                section_text += "\n"
            i += 4 + rec_len
        text += section_text
        text += "\n"
    text = re.sub('[^가-힣a-zA-Z0-9\s.,()-:]', '', text)
    text = text.replace('\x0b', '')
    
    return text

def hwpx_to_txt(file_path):
    text = ""
    with zipfile.ZipFile(file_path, 'r') as myzip:
        for file in myzip.namelist():
            if file.endswith('.xml'):
                with myzip.open(file) as myfile:
                    tree = parse(myfile)
                    root = tree.getroot()
                    for node in root.iter():
                        if node.tag.endswith('p'): # 'p' tag is for paragraphs
                            text += ''.join(node.itertext())
                            text += '\n' # add a new line at the end of each paragraph
        
        return text 
    
def jaccard_similarity_char_level(str1, str2):
    # 문자열에서 공백과 특수 문자를 제거합니다.
    str1 = re.sub(r"\W", "", str1)
    str2 = re.sub(r"\W", "", str2)
    
    # 문자열을 문자 단위로 분리하여 집합으로 변환합니다.
    set1 = set(str1)
    set2 = set(str2)
    
    # 교집합의 크기를 계산합니다.
    intersection = len(set1.intersection(set2))
    
    # 합집합의 크기를 계산합니다.
    union = len(set1.union(set2))
    
    # 교집합이 없으면 유사도는 0입니다.
    if union == 0:
        return 0
    
    # 자카드 유사도를 계산합니다.
    similarity = intersection / union
    return similarity

class MyApp(QWidget):

    def __init__(self):
        super().__init__()
        
        self.secondWindow = SecondWindow(self)  # SecondWindow 객체 생성
        self.initUI()

    def initUI(self):
        self.stackedWidget = QStackedWidget()

        self.firstPage = QWidget()
        self.secondPage = SecondWindow(self)

        self.firstPageLayout = QVBoxLayout()
        self.firstPage.setLayout(self.firstPageLayout)

        self.btn1 = QPushButton('개인정보 처리방침 불러오기', self)
        self.btn1.clicked.connect(self.openFile)
        self.firstPageLayout.addWidget(self.btn1)

        self.fileNameLabel = QLabel()
        self.firstPageLayout.addWidget(self.fileNameLabel)

        self.te1 = QTextEdit()
        self.firstPageLayout.addWidget(self.te1)

        self.stackedWidget.addWidget(self.firstPage)
        self.stackedWidget.addWidget(self.secondPage)

        self.mainLayout = QVBoxLayout()
        self.mainLayout.addWidget(self.stackedWidget)
        self.setLayout(self.mainLayout)

        self.btn2 = QPushButton('다음', self)
        self.btn2.clicked.connect(self.changePage)
        self.firstPageLayout.addWidget(self.btn2)

        self.setWindowTitle('Doc Doc')
        self.setWindowIcon(QIcon('docdoc.png'))
        self.setGeometry(1000, 1000, 1000, 1000)
        self.center()
        self.show()

    def center(self):
        # 창을 화면 가운데로 이동
        qr = self.frameGeometry()
        cp = QDesktopWidget().availableGeometry().center()
        qr.moveCenter(cp)
        self.move(qr.topLeft())

    def openFile(self):
        fname = QFileDialog.getOpenFileName(self, '개인정보 처리방침을 업로드해주세요')
 
        if fname[0]:
            extension = os.path.splitext(fname[0])[-1].lower()

            if extension == '.pdf':
                text = extract_text(fname[0])
            elif extension == '.hwp':
                text = get_hwp_text(fname[0])
            elif extension == '.hwpx':
                text = hwpx_to_txt(fname[0])
            elif extension == '.txt':
                with open(fname[0], "r", encoding="utf-8") as file:
                    text = file.read()
            elif extension in ['.docx']:
                text = read_docx(fname[0])
            elif extension in ['.pptx']:
                text = read_pptx(fname[0])
            elif extension in ['.png', '.jpg', '.jpeg']:
                text = ocr_image(fname[0])
            else:
                text = 'Unsupported file format.'

            self.secondWindow.file_texts.append(text)  # SecondWindow의 file_texts 리스트에 추가
            self.secondWindow.file_paths.append(fname[0])  # SecondWindow의 file_paths 리스트에 추가

            self.te1.setText(text)
            self.fileNameLabel.setText(os.path.basename(fname[0]))

    def changePage(self):
        self.stackedWidget.setCurrentWidget(self.secondPage)

class SecondWindow(QWidget):

    def __init__(self, parent):
        super().__init__(parent)
        self.main_app = parent
        self.file_open_counter = 0

         # 저장된 TfidfVectorizer 객체와 모델을 로드
        self.tfidf_vect = joblib.load('C:/Users/csm74/OneDrive/문서/카카오톡 받은 파일/굿/tfidf_vect.pkl')
        self.model = joblib.load('C:/Users/csm74/OneDrive/문서/카카오톡 받은 파일/굿/model.pkl')
        
        self.initUI()
        self.file_texts = []  # 파일 내용을 저장할 리스트1
        self.current_file_index = 0  # 현재 보여주고 있는 파일의 인덱스
        self.file_paths = []
    
    
    def initUI(self):
        self.vbox = QVBoxLayout()
        self.setLayout(self.vbox)

        self.btn_back = QPushButton('이전', self)
        self.btn_back.clicked.connect(self.goBack)

        self.hbox = QHBoxLayout()
        self.hbox.addWidget(self.btn_back)
        self.hbox.addStretch(1)  # 왼쪽 정렬을 위한 stretch 추가

        self.vbox.addLayout(self.hbox)  # hbox를 vbox에 추가

        self.btn1 = QPushButton('분석 파일 경로 선택', self)
        self.btn1.clicked.connect(self.openDirectory)
        self.vbox.addWidget(self.btn1)

        self.fileList = QTableWidget(self)
        self.fileList.setColumnCount(5)
        self.fileList.setHorizontalHeaderLabels(['Name', 'Creation Time', 'Size', 'Type', 'Path'])
        self.vbox.addWidget(self.fileList)

        self.btn_an = QPushButton('분석 시작', self)  # '분석' 버튼
        self.btn_an.clicked.connect(self.showAn)
        self.vbox.addWidget(self.btn_an)


    def goBack(self):
        self.main_app.stackedWidget.setCurrentWidget(self.main_app.firstPage)

    def openDirectory(self):
        dname = QFileDialog.getExistingDirectory(self, '파일 경로를 선택해주세요')

        if dname:
           # 새로운 디렉토리를 선택하면 파일 경로 리스트를 초기화
            self.file_paths = []
            self.show_files_in_folder(dname)

    def show_files_in_folder(self, folder_name):
        self.fileList.setRowCount(0)
        file_types = ['*.docx', '*.pptx', '*.xlsx', '*.hwp', '*.txt', '*.png', '*.jpg', '*.pdf','*.hwpx']
        for file_type in file_types:
            for file_name in glob.glob(os.path.join(folder_name, file_type)):
                self.file_paths.append(file_name)
                creation_time = time.strftime('%Y-%m-%d %H:%M:%S', time.localtime(os.path.getctime(file_name)))
                base_name = os.path.basename(file_name)
                file_size = os.path.getsize(file_name) / 1024
                file_type = os.path.splitext(file_name)[-1]
                self.fileList.insertRow(self.fileList.rowCount())
                self.fileList.setItem(self.fileList.rowCount()-1, 0, QTableWidgetItem(base_name))
                self.fileList.setItem(self.fileList.rowCount()-1, 1, QTableWidgetItem(creation_time))
                self.fileList.setItem(self.fileList.rowCount()-1, 2, QTableWidgetItem(f'{file_size:.2f} KB'))
                self.fileList.setItem(self.fileList.rowCount()-1, 3, QTableWidgetItem(file_type))
                self.fileList.setItem(self.fileList.rowCount()-1, 4, QTableWidgetItem(file_name))


    def showNextFile(self):
        self.current_file_index += 1  # 인덱스를 1 증가
        if self.current_file_index >= len(self.file_paths):  # 리스트의 끝에 도달했을 경우 인덱스를 0으로 초기화
            self.current_file_index = 0
        if self.dialog:  # 만약 dialog가 존재하면
            file_path = self.file_paths[self.current_file_index]
            text = self.getFileContent(file_path)
            self.dialog.te_current.setText(text)  # dialog의 QTextEdit에 다음 파일 내용을 표시
            file_name = os.path.basename(file_path)  # 파일 경로에서 파일 이름 추출
            self.dialog.setWindowTitle(file_name)
            predictions = self.predict_topic(text)
            self.dialog.te_current.setText("예측 결과:") 


            # Loop through predictions to set background colors
            for sent, pred in predictions:
                if "분류: 1.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFCCCC"
                elif "분류: 2.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFFFCC"
                elif "분류: 3.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCFFCC"
                elif "분류: 4.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCE5FF"
                else:
                    background_color = ""  # Set a default color if not matched

                # Set background color for each line
                self.dialog.te_current.append(f'<span style="background-color: {background_color};">{sent} | {pred}</span>')
       
    def showAn(self):
        if self.file_paths:
            file_path = self.file_paths[self.current_file_index]
            text = self.getFileContent(file_path)
            file_name = os.path.basename(file_path)
            self.dialog = SecondWindow.FileDisplayDialog(text, file_name, self)
            self.dialog.setParentShowMethod(self.showAn2)
            self.dialog.show()

            predictions = self.predict_topic(text)
            self.dialog.te_current.setText("예측 결과:")

            # Loop through predictions to set background colors
            for sent, pred in predictions:
                if "분류: 1.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFCCCC"
                elif "분류: 2.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFFFCC"
                elif "분류: 3.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCFFCC"
                elif "분류: 4.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCE5FF"
                else:
                    background_color = ""  # Set a default color if not matched

                # Set background color for each line
                self.dialog.te_current.append(f'<span style="background-color: {background_color};">{sent} | {pred}</span>')


    def showAn2(self):
        self.current_file_index += 1  # 인덱스를 1 증가
        if self.current_file_index >= len(self.file_paths):  # 리스트의 끝에 도달했을 경우 인덱스를 0으로 초기화
            self.current_file_index = 0
        if self.dialog:  # 만약 dialog가 존재하면
            file_path = self.file_paths[self.current_file_index]
            text = self.getFileContent(file_path)
            #self.dialog.te.setText(text)  # dialog의 QTextEdit에 다음 파일 내용을 표시
            file_name = os.path.basename(file_path)
            self.dialog.setWindowTitle(file_name)

            predictions = self.predict_topic(text)
            self.dialog.te_current.setText("예측 결과:")

            # Loop through predictions to set background colors
            for sent, pred in predictions:
                if "분류: 1.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFCCCC"
                elif "분류: 2.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#FFFFCC"
                elif "분류: 3.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCFFCC"
                elif "분류: 4.0" in f"문장: {sent}, 분류: {pred}":
                    background_color = "#CCE5FF"
                else:
                    background_color = ""  # Set a default color if not matched

                # Set background color for each line
                self.dialog.te_current.append(f'<span style="background-color: {background_color};">{sent} | {pred}</span>')
        
    

    def getFileContent(self, file_path):
        extension = os.path.splitext(file_path)[-1].lower()

        if extension == '.pdf':
            text = extract_text(file_path)
        elif extension == '.hwp':
            text = get_hwp_text(file_path)
        elif extension == '.hwpx':
            text = hwpx_to_txt(file_path)
        elif extension == '.txt':
            with open(file_path, "r", encoding="utf-8") as file:  # 인코딩을 'utf-8'로 설정
                text = file.read()
        elif extension in ['.docx']:
            text = read_docx(file_path)
        elif extension in ['.pptx']:
            text = read_pptx(file_path)
        elif extension in ['.png', '.jpg', '.jpeg']:
            text = ocr_image(file_path)
        else:
            text = 'Unsupported file format.'

        return text
    
    def preprocess(self, text):
        okt = Okt()
        def okt_clean(text):
            clean_text = []
            for word in okt.pos(text, stem=True):
                if word[1] not in ['Josa', 'Eomi','Punctuation']:
                    clean_text.append(word[0])
            return " ".join(clean_text)

        stops = ['합니다','하는','할','하고','한다', '그리고', '입니다', '그','등','이런','것','및','제','더','있습니다.','관련','대한','?']
        tokens = text.split(' ')
        meaningful_words = [w for w in tokens if not w in stops]
        text = ' '.join(meaningful_words)

        text = re.sub('[^가-힣ㄱ-ㅎㅏ-ㅣa-zA-Z]', ' ', text)
        text = re.sub('[\s]+', ' ', text)
        text = text.lower()

        return text

    def predict_topic(self, text): 
        sentences = text.split('\n')  
        self.pass_sentences = ["<개인정보 제공 및 이용 동의서> ", "<개인정보 수집 및 이용 동의서> ", "1. 개인정보 수집 및 이용 목적", "2. 개인정보 수집 항목", "3. 개인정보의 보유 및 이용 기간", "4. 동의 거부 및 동의 거부 시 불이익 내용","2) 정보 제공"," <개인정보 제공 및 이용 동의서> ","<개인정보 제공 및 이용 동의서> ","개인정보 제공 및 이용 동의서 "]
        self.pass_words = ["다음과", "(서명)", "같이", "1.0"]
        self.one_words = ["환불", "기업", "가입형", "표준성과한마당,", "자동입력", "수수료", "정보제공", "회원가입", "해결", "채용성공수수료"]
        self.two_words = ["닉네임", "전공", "세부전공", "이메일", "제목", "연령", "관심분야", "장애여부,", "보훈여부,", "학력사항", "경력사항", "소속명", "일반전화번호", "아이디", "비밀번호", "필수 :", "선택 :","필수:","공개여부","휴대전화번호", "필수", "선택"] 
        self.three_words = ["10년동안","회신일로부터","12개월까지","보유기간은","보관됩니다.", "폐기", "법률에"]
        predictions = []

        # 문장 토큰화 함수 정의
        def sentence_tokenize(sentence):
            return sentence.strip().replace('\t', '').split()
        

        for sentence in sentences:
            sentence_tokens = sentence_tokenize(sentence)
            pass_tokens = [sentence_tokenize(pass_sentence) for pass_sentence in self.pass_sentences]
            pass_flag1 = any(word in sentence_tokens for word in self.pass_words)
            pass_flag2 = any(sentence_tokens == pass_sentence for pass_sentence in pass_tokens)
            one_flag = any(word in sentence_tokens for word in self.one_words)
            two_flag = any(word in sentence_tokens for word in self.two_words)
            three_flag = any(word in sentence_tokens for word in self.three_words)

            if pass_flag1:
                predictions.append((sentence.strip(), 'pass'))
            elif three_flag:
                predictions.append((sentence.strip(), '3.0'))   
            elif one_flag:
                predictions.append((sentence.strip(), '1.0'))  
            elif pass_flag2:
                predictions.append((sentence.strip(), 'pass'))
            elif two_flag:
                predictions.append((sentence.strip(), '2.0'))
            else:
                preprocessed = self.preprocess(sentence)
                text_feature_tfidf = self.tfidf_vect.transform([preprocessed])
                probas = self.model.predict_proba(text_feature_tfidf)
                max_proba = max(probas[0])
                max_idx = list(probas[0]).index(max_proba)
                pred = self.model.classes_[max_idx]

                if max_proba >= 0.6:
                    predictions.append((sentence.strip(), pred))
                else:
                    predictions.append((sentence.strip(), 'pass'))

        return predictions

        
    class FileDisplayDialog(QDialog):
        def __init__(self, text, file_name, parent=None):
            super().__init__(parent)
            self.setWindowTitle(file_name)
            self.setGeometry(1000, 1000, 1000, 1000)

            self.main_layout = QVBoxLayout()  # 메인 레이아웃을 생성합니다.
            self.layout = QHBoxLayout()  # 처음에 받아온 문서와 현재 보고 있는 문서를 보여주는 레이아웃을 생성합니다.

        # 처음에 받아온 문서를 보여줄 QTextEdit
            
            initial_policy_text = parent.main_app.te1.toPlainText()  # 초기 내용 저장
            self.te_policy = QTextEdit()
            self.te_policy.setText(initial_policy_text)
            self.initial_policy_text = initial_policy_text  # self에 저장하여 객체 내에서 사용할 수 있도록 함
            self.layout.addWidget(self.te_policy)
            self.te_policy.hide() 

        # 지금 보고 있는 문서를 보여줄 QTextEdit
            self.te_current = QTextEdit()
            self.te_current.setText(text)
            self.layout.addWidget(self.te_current)

            self.bottom_text = QLabel()
            self.bottom_text.setText(
            "<html>개인정보 보호법 (제 15조 2항)<br>"
            "② 개인정보처리자는 제1항제1호에 따른 동의를 받을 때에는 다음 각 호의 사항을 정보주체에게 알려야 한다.<br>" 
            "다음 각 호의 어느 하나의 사항을 변경하는 경우에도 이를 알리고 동의를 받아야 한다.<br>"
            "1. 개인정보의 수집ㆍ이용 목적<br>"
            "2. 수집하려는 개인정보의 항목<br>"
            "3. 개인정보의 보유 및 이용 기간<br>"
            "4. 동의를 거부할 권리가 있다는 사실 및 동의 거부에 따른 불이익이 있는 경우에는 그 불이익의 내용<br>"
            "<br>문장분류: <br><br>"
            "<span style='background-color: #FFCCCC;'>1.0 : 개인정보의 수집ㆍ이용 목적</span><br><br>"
            "<span style='background-color: #FFFFCC;'>2.0 : 수집하려는 개인정보의 항목</span><br><br>"
            "<span style='background-color: #CCFFCC;'>3.0 : 개인정보의 보유 및 이용 기간</span><br><br>"
            "<span style='background-color: #CCE5FF;'>4.0 : 동의를 거부할 권리가 있다는 사실 및 동의 거부에 따른 불이익이 있는 경우에는 그 불이익의 내용</span><br><br></html>"
        )
            self.main_layout.addWidget(self.bottom_text)
            self.setLayout(self.main_layout)

        # '누락 사항 분석' 버튼
            self.btn_policy = QPushButton('누락 사항 분석', self)
            self.btn_policy.clicked.connect(self.showBothDocuments)

        # '다음 파일 확인' 버튼
            self.btn_next = QPushButton('다음 파일 확인', self)
            self.btn_next.clicked.connect(parent.showNextFile)
            
        # '이전' 버튼
            self.btn_prev = QPushButton('이전', self)
            self.btn_prev.clicked.connect(self.goBack)
            self.btn_prev.hide()  # 처음에는 숨겨둡니다.
            self.resetPolicyText()

            self.main_layout.addLayout(self.layout)
            self.main_layout.addWidget(self.btn_policy)
            self.main_layout.addWidget(self.btn_next)
            self.main_layout.addWidget(self.btn_prev)

            self.setLayout(self.main_layout)

        def resetPolicyText(self):
            self.te_policy.setText(self.initial_policy_text) 

        def check_missing_sentences(self, text):
            # 1.0 ~ 4.0 에 해당하는 문장이 있는지 체크하고, 없으면 경고 메시지를 반환
            missing_sentences = []
            for i in range(1, 5):
                if f"{i}.0" not in text:
                    missing_sentences.append(f"경고 : 문장 분류 {i}.0에 해당하는 문장이 없습니다.\n문서를 다시 검토해주세요.")
            return missing_sentences

        def showBothDocuments(self):

            # 왼쪽 텍스트 박스에 누락 사항이 발견된 문서의 1.0 문장과 주변 3문장을 표시
            matching_sentences = self.find_matching_sentences(self.te_current.toPlainText(), self.te_policy.toPlainText())
            self.te_policy.setText('\n'.join(matching_sentences))

            # '처리방침 확인하기' 버튼을 누르면 처음에 받아온 문서와 지금 보고 있는 문서를 동시에 보여줍니다.
            self.te_policy.show()
            self.te_current.show()
            self.btn_policy.hide()
            self.btn_next.hide()
            self.btn_prev.show()
           

            self.btn_compare = QPushButton('처리방침과 비교하기', self)
            self.btn_compare.clicked.connect(self.compareWithPolicy)
            self.main_layout.addWidget(self.btn_compare)

            # 추가된 부분: 경고 메시지 확인
            missing_sentences = self.check_missing_sentences(self.te_current.toPlainText())
            if missing_sentences:
                warning_message = "\n".join(missing_sentences)
                QMessageBox.warning(self, "경고", warning_message)
            else:
                # 1.0 ~ 4.0 에 해당하는 문장이 모두 있으면 "누락 사항이 발견되지 않았습니다." 메시지를 띄움
                QMessageBox.information(self, "알림", "누락 사항이 발견되지 않았습니다.")


        def find_matching_sentences(self, current_text, policy_text):
    # 현재 문서의 1.0과 3.0에 해당하는 문장 찾기
                current_sentences = current_text.split('\n')
                matching_sentence_1_0 = next((sentence.strip() for sentence in current_sentences if "1.0" in sentence), None)
                matching_sentence_3_0 = next((sentence.strip() for sentence in current_sentences if "3.0" in sentence), None)

                matching_policy_contexts = []
                self.resetPolicyText()

                if matching_sentence_1_0 is not None or matching_sentence_3_0 is not None:
        # 개인정보 처리방침 파일에서 유사한 문장 찾기
                    policy_sentences = policy_text.split('\n')
        
                    keywords = ['수집정보', '항목', '기간', '수집항목', '목적']  # 검사하려는 단어들을 리스트로 만듭니다.

                    for matching_sentence in [matching_sentence_1_0, matching_sentence_3_0]:
                        if matching_sentence is not None:
                            for i, policy_sentence in enumerate(policy_sentences):
                                if jaccard_similarity_char_level(matching_sentence, policy_sentence.strip()) >= 0.5:
                                    start_index = max(0, i - 5)
                                    end_index = min(len(policy_sentences), i + 5)
                                    context_sentences = policy_sentences[start_index:end_index]

                        # 'keywords' 리스트에 있는 단어 중 하나라도 문장에 포함되어 있는지 확인
                                    if any(keyword in sentence for sentence in context_sentences for keyword in keywords):
                                        print(f"유사한 문장: {policy_sentence}")
                                        matching_policy_context = '\n'.join(context_sentences)
                                        matching_policy_contexts.append(matching_policy_context)

                if matching_policy_contexts:
                    return matching_policy_contexts

                return ["해당되는 개인정보 처리방침 문장을 찾을 수 없습니다."]

                              
        def compareWithPolicy(self):
    # document_text와 policy_text를 가져옵니다.
            
            document_text = self.te_current.toPlainText()
            policy_text = self.te_policy.toPlainText()
            
    # document_text에서 2.0과 3.0이 포함된 문장을 찾습니다.
            document_target_sentences = []
            document_sentences = document_text.split('\n')
            for sentence in document_sentences:
                if "2.0" in sentence:
                    document_target_sentences.append(sentence.strip())
                if "3.0" in sentence:
                    document_target_sentences.append(sentence.strip())
            

            policy_sentences = policy_text.split('\n')
            similar_sentences = {}
            for document_sentence in document_target_sentences:
                for policy_sentence in policy_sentences:
                    similarity = jaccard_similarity_char_level(document_sentence, policy_sentence)
                    if similarity > 0.4564:  # 유사도 임계값 조정 가능
                        if document_sentence not in similar_sentences:
                            similar_sentences[document_sentence] = []
                        similar_sentences[document_sentence].append((policy_sentence, similarity))

# 유사한 문장과 그 유사도를 출력합니다.
            message = "<br/>"
            for document_sentence, similar_sentence_info in similar_sentences.items():
                message += f"<font color='#000000'>동의서 문장:<br/> {document_sentence}</font><br/>"
                for info in similar_sentence_info:
                    similarity_percentage = round(info[1] * 100, 2)

                    if similarity_percentage <= 50:
                        background_color = "#ED9595"
                        label = "(위험)"
                    elif 50 < similarity_percentage <= 80:
                        background_color = "#FFE08C"
                        label = "(경고)"
                    else:
                        background_color = "#CEF279"
                        label = "(안전)"

                    similarity_text = f"유사도: {similarity_percentage}% {label}"
                    similarity_text_with_background = f"<font style='background-color: {background_color}'>{similarity_text}</font>"

                    message += f"<font color='#013ADF'>개인정보처리방침 문장: <br/>{info[0]}<br/> <font color='#000000'>{similarity_text_with_background}<br/>"
                message += "<br/>"

            msg = QMessageBox()
            msg.setWindowTitle("유사 문장 비교 결과")
            msg.setText(message)
            msg.exec_()




            
        def goBack(self):
            self.te_policy.hide()
            self.te_current.show()
            self.btn_policy.show()
            self.btn_next.show()
            self.btn_prev.hide()
            self.btn_compare.hide()
            self.resetPolicyText()
        
        def setParentShowMethod(self, method):
            # 부모(parent)의 메서드를 동적으로 변경하는 메서드
            self.parent_show_method = method

        def call_parent_method(self):
            # 저장된 부모(parent)의 메서드 호출
            if self.parent_show_method:
                self.parent_show_method()


if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = MyApp()
    sys.exit(app.exec_())
